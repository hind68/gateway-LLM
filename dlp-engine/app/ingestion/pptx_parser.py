from pptx import Presentation

from app.ingestion.ooxml_media import extract_text_from_embedded_images


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)

        # has_notes_slide is checked first because the notes_slide
        # getter itself creates an (empty) notes slide as a side effect
        # if one doesn't already exist - this is a read-only extraction,
        # so it shouldn't be modifying the file's structure as it goes.
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                parts.append(notes_text)

    # A pasted screenshot or scanned photo (e.g. of an ID card) on a
    # slide has text too, same reasoning as pdf_parser.py's OCR pass.
    parts.extend(extract_text_from_embedded_images(path, "ppt/media/"))

    return "\n".join(parts)
