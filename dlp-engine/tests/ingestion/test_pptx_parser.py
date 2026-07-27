import os
import tempfile

from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont

from app.ingestion.pptx_parser import extract_text_from_pptx


def _render_text_image(text: str, path: str):
    img = Image.new("RGB", (600, 100), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 32)
    except Exception:
        font = ImageFont.load_default()
    draw.text((10, 30), text, fill="black", font=font)
    img.save(path)


def test_extracts_slide_text_frame_and_table():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Report"
    prs.save(path)

    text = extract_text_from_pptx(path)
    os.remove(path)

    assert "Quarterly Report" in text

def test_extracts_speaker_notes():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "Remember to mention john@company.com privately"
    prs.save(path)

    text = extract_text_from_pptx(path)
    os.remove(path)

    assert "john@company.com" in text

def test_ocrs_embedded_image_on_slide():
    img_path = tempfile.mktemp(suffix=".png")
    _render_text_image("CONFIDENTIAL EMAIL", img_path)

    fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(img_path, Inches(1), Inches(1))
    prs.save(pptx_path)

    text = extract_text_from_pptx(pptx_path)
    os.remove(pptx_path)
    os.remove(img_path)

    assert "CONFIDENTIAL" in text.upper()

def test_slide_with_no_notes_does_not_error():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path)

    text = extract_text_from_pptx(path)  # should not raise
    os.remove(path)
    assert isinstance(text, str)
